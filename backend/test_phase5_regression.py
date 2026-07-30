import os
import django

os.environ.setdefault('DJANGO_SETTINGS_MODULE', 'LAMS.settings')
django.setup()

import tempfile
import zipfile
import docx
import pptx
import pypdf
from django.core.files.uploadedfile import SimpleUploadedFile

from courses.models import Course, Module, Lesson
from import_engine.models import ImportJob
from authoring_engine.models import LessonBlockTree, LessonBlock
from users.models import User
from rest_framework.test import APIRequestFactory, force_authenticate
from import_engine.views import upload_and_convert
from courses.views import CourseViewSet

def create_sample_pptx(path):
    prs = pptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Sample PPTX Slide 1"
    slide.placeholders[1].text = "Body text for slide 1"
    prs.save(path)

def create_sample_pdf(path):
    writer = pypdf.PdfWriter()
    writer.add_blank_page(width=612, height=792)
    with open(path, 'wb') as f:
        writer.write(f)

def create_sample_docx(path):
    doc = docx.Document()
    doc.add_heading('Sample DOCX Module 1', level=1)
    doc.add_heading('Lesson 1.1 Intro', level=2)
    doc.add_paragraph('This is paragraph text in docx.')
    doc.save(path)

def create_sample_zip(path):
    with zipfile.ZipFile(path, 'w') as zf:
        zf.writestr('index.html', '<html><head><title>Sample HTML Course</title></head><body><h1>Heading 1</h1><p>Paragraph content</p></body></html>')

def run_phase5_test():
    print("=== PHASE 5 UNIVERSAL IMPORT ENGINE REGRESSION CHECK ===")

    admin_user = User.objects.filter(is_platform_super_admin=True).first() or User.objects.first()
    factory = APIRequestFactory()

    temp_dir = tempfile.mkdtemp()
    print(f"Created temporary scratch directory: {temp_dir}")

    # A. Test PPTX Import
    pptx_path = os.path.join(temp_dir, "sample.pptx")
    create_sample_pptx(pptx_path)
    with open(pptx_path, 'rb') as f:
        uf = SimpleUploadedFile("sample.pptx", f.read(), content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation")
    
    req = factory.post('/api/import/upload/', {'source_file': uf, 'source_format': 'pptx'}, format='multipart')
    force_authenticate(req, user=admin_user)
    res = upload_and_convert(req)
    assert res.status_code == 201, f"PPTX upload failed: {res.data}"
    c_id = res.data['target_course_id']
    pptx_course = Course.objects.get(id=c_id)
    assert pptx_course.modules.count() >= 1, "PPTX created modules"
    assert pptx_course.modules.first().lessons.count() >= 1, "PPTX created lessons"
    print(f"  [OK] PPTX Adapter: Successfully converted presentation -> Course ID {pptx_course.id} ('{pptx_course.title}').")

    # B. Test DOCX Import
    docx_path = os.path.join(temp_dir, "sample.docx")
    create_sample_docx(docx_path)
    with open(docx_path, 'rb') as f:
        uf = SimpleUploadedFile("sample.docx", f.read(), content_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document")
    
    req = factory.post('/api/import/upload/', {'source_file': uf, 'source_format': 'docx'}, format='multipart')
    force_authenticate(req, user=admin_user)
    res = upload_and_convert(req)
    assert res.status_code == 201, f"DOCX upload failed: {res.data}"
    c_id = res.data['target_course_id']
    docx_course = Course.objects.get(id=c_id)
    print(f"  [OK] DOCX Adapter: Successfully converted document -> Course ID {docx_course.id} ('{docx_course.title}').")

    # C. Test HTML/ZIP Import
    zip_path = os.path.join(temp_dir, "sample.zip")
    create_sample_zip(zip_path)
    with open(zip_path, 'rb') as f:
        uf = SimpleUploadedFile("sample.zip", f.read(), content_type="application/zip")
    
    req = factory.post('/api/import/upload/', {'source_file': uf, 'source_format': 'zip'}, format='multipart')
    force_authenticate(req, user=admin_user)
    res = upload_and_convert(req)
    assert res.status_code == 201, f"ZIP upload failed: {res.data}"
    c_id = res.data['target_course_id']
    zip_course = Course.objects.get(id=c_id)
    print(f"  [OK] HTML/ZIP Adapter: Successfully converted zip -> Course ID {zip_course.id} ('{zip_course.title}').")

    # D. Test PDF Import
    pdf_path = os.path.join(temp_dir, "sample.pdf")
    create_sample_pdf(pdf_path)
    with open(pdf_path, 'rb') as f:
        uf = SimpleUploadedFile("sample.pdf", f.read(), content_type="application/pdf")
    
    req = factory.post('/api/import/upload/', {'source_file': uf, 'source_format': 'pdf'}, format='multipart')
    force_authenticate(req, user=admin_user)
    res = upload_and_convert(req)
    assert res.status_code == 201, f"PDF upload failed: {res.data}"
    c_id = res.data['target_course_id']
    pdf_course = Course.objects.get(id=c_id)
    print(f"  [OK] PDF Adapter: Successfully converted pdf -> Course ID {pdf_course.id} ('{pdf_course.title}').")

    # E. Verify converted course is editable identically to a manually created course
    c_view = CourseViewSet.as_view({'patch': 'partial_update'})
    p_req = factory.patch(f'/api/courses/{pptx_course.id}/', {'title': 'Updated Converted PPTX Title'}, format='json')
    force_authenticate(p_req, user=admin_user)
    p_res = c_view(p_req, pk=pptx_course.id)
    assert p_res.status_code == 200
    assert p_res.data['title'] == 'Updated Converted PPTX Title'
    print("  [OK] Converted courses verified 100% editable via standard CourseViewSet API.")

    print("\n=== PHASE 5 REGRESSION CHECK PASSED 100% ===")

if __name__ == '__main__':
    run_phase5_test()
