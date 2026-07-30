import os
import re
import uuid
import hashlib
import zipfile
import tempfile
from io import BytesIO
from django.db import transaction
from django.core.files.base import ContentFile

import pptx
import pypdf
import docx
from bs4 import BeautifulSoup
import mutagen
import pdf2image

from authoring_engine.models import (
    AuthoringAsset, LessonBlockTree, LessonBlock,
    ReadingContent, InteractionBlock, KCQuestion, ScenarioNode
)
from courses.models import Course, Module, Lesson
from courses.scorm_engine import parse_imsmanifest

class ImportAdapterException(Exception):
    pass

class BaseAdapter:
    def __init__(self, file_path, organization, user, target_course=None):
        self.file_path = file_path
        self.organization = organization
        self.user = user
        self.target_course = target_course
        self.warnings = []
        self.errors = []

    def parse(self):
        """Returns intermediate AST: { 'title': '...', 'modules': [ { 'title': '...', 'lessons': [ ... ] } ] }"""
        raise NotImplementedError


class PPTXAdapter(BaseAdapter):
    def parse(self):
        try:
            prs = pptx.Presentation(self.file_path)
        except Exception as e:
            raise ImportAdapterException(f"Failed to open PPTX file: {str(e)}")

        total_slides = len(prs.slides)
        if total_slides == 0:
            raise ImportAdapterException("PPTX presentation contains no slides.")

        failed_slides = 0
        lessons = []

        for idx, slide in enumerate(prs.slides, start=1):
            try:
                slide_title = f"Slide {idx}"
                blocks = []
                
                # Search for slide title shape
                for shape in slide.shapes:
                    if shape.has_text_frame and shape.text.strip():
                        if shape == slide.shapes.title or not blocks:
                            slide_title = shape.text.strip().split('\n')[0]
                            blocks.append({
                                'type': 'heading',
                                'payload': {'html': f"<h2>{slide_title}</h2>", 'markdown': slide_title}
                            })
                        else:
                            text_content = shape.text.strip()
                            if text_content:
                                blocks.append({
                                    'type': 'paragraph',
                                    'payload': {'html': f"<p>{text_content}</p>", 'markdown': text_content}
                                })

                if not blocks:
                    blocks.append({
                        'type': 'paragraph',
                        'payload': {'html': f"<p>Slide {idx} Content</p>", 'markdown': f"Slide {idx} Content"}
                    })

                lessons.append({
                    'title': slide_title,
                    'type': 'reading',
                    'duration': '5 min',
                    'blocks': blocks
                })

            except Exception as slide_err:
                failed_slides += 1
                self.warnings.append(f"Slide {idx} parsing warning: {str(slide_err)}")
                lessons.append({
                    'title': f"Slide {idx} (Conversion Warning)",
                    'type': 'reading',
                    'duration': '5 min',
                    'blocks': [{
                        'type': 'callout',
                        'payload': {
                            'html': f"<p><strong>Warning:</strong> Content on Slide {idx} could not be automatically converted.</p>",
                            'markdown': f"Warning: Content on Slide {idx} could not be automatically converted."
                        }
                    }]
                })

        # Check section 4.3 rollback threshold (>30% failure)
        failure_rate = (failed_slides / total_slides) * 100
        if failure_rate > 30.0:
            raise ImportAdapterException(f"Import failed: {failure_rate:.1f}% of slides failed parsing (>30% threshold).")

        main_title = os.path.basename(self.file_path).replace('.pptx', '').replace('_', ' ').title()
        return {
            'title': main_title,
            'modules': [
                {
                    'title': 'Module 1: Presentation Slides',
                    'summary': f'Imported from {os.path.basename(self.file_path)}',
                    'lessons': lessons
                }
            ]
        }


class PDFAdapter(BaseAdapter):
    def parse(self):
        try:
            reader = pypdf.PdfReader(self.file_path)
        except Exception as e:
            raise ImportAdapterException(f"Failed to open PDF file: {str(e)}")

        total_pages = len(reader.pages)
        if total_pages == 0:
            raise ImportAdapterException("PDF document contains no pages.")

        failed_pages = 0
        lessons = []

        for idx, page in enumerate(reader.pages, start=1):
            try:
                text = page.extract_text() or ''
                text = text.strip()

                lines = [line.strip() for line in text.split('\n') if line.strip()]
                page_title = lines[0] if lines else f"Page {idx}"
                body_text = "\n\n".join(lines[1:]) if len(lines) > 1 else text

                blocks = [
                    {'type': 'heading', 'payload': {'html': f"<h2>{page_title}</h2>", 'markdown': page_title}},
                    {'type': 'paragraph', 'payload': {'html': f"<p>{body_text or 'Document Page'}</p>", 'markdown': body_text or 'Document Page'}}
                ]

                lessons.append({
                    'title': page_title[:200],
                    'type': 'reading',
                    'duration': '5 min',
                    'blocks': blocks
                })

            except Exception as page_err:
                failed_pages += 1
                self.warnings.append(f"Page {idx} warning: {str(page_err)}")
                lessons.append({
                    'title': f"Page {idx} (Warning)",
                    'type': 'reading',
                    'duration': '5 min',
                    'blocks': [{
                        'type': 'callout',
                        'payload': {'html': f"<p>Warning: Page {idx} content could not be parsed.</p>", 'markdown': f"Warning: Page {idx} content could not be parsed."}
                    }]
                })

        failure_rate = (failed_pages / total_pages) * 100
        if failure_rate > 30.0:
            raise ImportAdapterException(f"Import failed: {failure_rate:.1f}% of pages failed parsing (>30% threshold).")

        main_title = os.path.basename(self.file_path).replace('.pdf', '').replace('_', ' ').title()
        return {
            'title': main_title,
            'modules': [
                {
                    'title': 'Module 1: Document Chapters',
                    'summary': f'Imported from {os.path.basename(self.file_path)}',
                    'lessons': lessons
                }
            ]
        }


class DOCXAdapter(BaseAdapter):
    def parse(self):
        try:
            doc = docx.Document(self.file_path)
        except Exception as e:
            raise ImportAdapterException(f"Failed to open DOCX document: {str(e)}")

        modules = []
        current_module = {'title': 'Module 1: General Content', 'summary': 'Docx import', 'lessons': []}
        current_lesson = {'title': 'Introduction', 'type': 'reading', 'duration': '5 min', 'blocks': []}

        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue

            style_name = para.style.name.lower() if para.style else ''
            if 'heading 1' in style_name:
                if current_lesson['blocks']:
                    current_module['lessons'].append(current_lesson)
                    current_lesson = {'title': text, 'type': 'reading', 'duration': '5 min', 'blocks': []}
                if current_module['lessons']:
                    modules.append(current_module)
                    current_module = {'title': text, 'summary': '', 'lessons': []}
            elif 'heading 2' in style_name:
                if current_lesson['blocks']:
                    current_module['lessons'].append(current_lesson)
                current_lesson = {'title': text, 'type': 'reading', 'duration': '5 min', 'blocks': []}
            else:
                current_lesson['blocks'].append({
                    'type': 'paragraph',
                    'payload': {'html': f"<p>{text}</p>", 'markdown': text}
                })

        if current_lesson['blocks'] or not current_module['lessons']:
            if not current_lesson['blocks']:
                current_lesson['blocks'].append({'type': 'paragraph', 'payload': {'html': '<p>Content</p>', 'markdown': 'Content'}})
            current_module['lessons'].append(current_lesson)
        if current_module['lessons']:
            modules.append(current_module)

        main_title = os.path.basename(self.file_path).replace('.docx', '').replace('_', ' ').title()
        return {
            'title': main_title,
            'modules': modules
        }


class VideoAdapter(BaseAdapter):
    def parse(self):
        rel_url = f"/media/import_scratch/{os.path.basename(self.file_path)}"
        main_title = os.path.basename(self.file_path).rsplit('.', 1)[0].replace('_', ' ').title()

        return {
            'title': main_title,
            'modules': [
                {
                    'title': 'Module 1: Video Content',
                    'summary': 'Imported video lesson',
                    'lessons': [
                        {
                            'title': main_title,
                            'type': 'video',
                            'duration': '10 min',
                            'video_url': rel_url,
                            'blocks': []
                        }
                    ]
                }
            ]
        }


class AudioAdapter(BaseAdapter):
    def parse(self):
        rel_url = f"/media/import_scratch/{os.path.basename(self.file_path)}"
        main_title = os.path.basename(self.file_path).rsplit('.', 1)[0].replace('_', ' ').title()

        return {
            'title': main_title,
            'modules': [
                {
                    'title': 'Module 1: Audio Content',
                    'summary': 'Imported audio lesson',
                    'lessons': [
                        {
                            'title': main_title,
                            'type': 'reading',
                            'duration': '5 min',
                            'blocks': [
                                {
                                    'type': 'audio',
                                    'payload': {
                                        'html': f'<audio controls src="{rel_url}"></audio>',
                                        'markdown': f'[Audio Track]({rel_url})'
                                    }
                                }
                            ]
                        }
                    ]
                }
            ]
        }


class HTMLZipAdapter(BaseAdapter):
    def parse(self):
        if not zipfile.is_zipfile(self.file_path):
            raise ImportAdapterException("File is not a valid ZIP archive.")

        with zipfile.ZipFile(self.file_path, 'r') as zf:
            namelist = zf.namelist()
            html_files = [f for f in namelist if f.endswith('.html') or f.endswith('.htm')]
            if not html_files:
                raise ImportAdapterException("ZIP archive contains no HTML files.")

            index_file = 'index.html' if 'index.html' in html_files else html_files[0]
            with zf.open(index_file) as f:
                soup = BeautifulSoup(f.read().decode('utf-8', errors='ignore'), 'html.parser')

            title = soup.title.string.strip() if soup.title and soup.title.string else 'HTML Package'
            blocks = []
            for tag in soup.find_all(['h1', 'h2', 'h3', 'p', 'table']):
                if tag.name.startswith('h'):
                    blocks.append({'type': 'heading', 'payload': {'html': str(tag), 'markdown': tag.text.strip()}})
                elif tag.name == 'p':
                    blocks.append({'type': 'paragraph', 'payload': {'html': str(tag), 'markdown': tag.text.strip()}})
                elif tag.name == 'table':
                    blocks.append({'type': 'table', 'payload': {'html': str(tag), 'markdown': tag.text.strip()}})

            if not blocks:
                blocks.append({'type': 'paragraph', 'payload': {'html': '<p>HTML Content</p>', 'markdown': 'HTML Content'}})

            return {
                'title': title,
                'modules': [
                    {
                        'title': 'Module 1: Web Content',
                        'summary': f'Extracted from {index_file}',
                        'lessons': [
                            {'title': title, 'type': 'reading', 'duration': '5 min', 'blocks': blocks}
                        ]
                    }
                ]
            }


class SCORMAdapter(BaseAdapter):
    def parse(self):
        target_xml = self.file_path

        # If file is a ZIP archive, extract to temporary directory
        if zipfile.is_zipfile(self.file_path):
            temp_dir = tempfile.mkdtemp()
            with zipfile.ZipFile(self.file_path, 'r') as zf:
                zf.extractall(temp_dir)
            target_xml = os.path.join(temp_dir, 'imsmanifest.xml')

        manifest_data = parse_imsmanifest(target_xml)
        title = manifest_data.get('title') or 'SCORM Package'
        sco_items = manifest_data.get('sco_items', [])

        lessons = []
        for idx, item in enumerate(sco_items, start=1):
            lessons.append({
                'title': item.get('title') or f"SCO Item {idx}",
                'type': 'reading',
                'duration': '10 min',
                'blocks': [
                    {
                        'type': 'callout',
                        'payload': {
                            'html': f"<p><strong>SCORM SCO:</strong> {item.get('title')} (Launch: {item.get('launch_href')})</p>",
                            'markdown': f"SCORM SCO: {item.get('title')} ({item.get('launch_href')})"
                        }
                    }
                ]
            })

        if not lessons:
            lessons.append({
                'title': title,
                'type': 'reading',
                'duration': '10 min',
                'blocks': [{'type': 'paragraph', 'payload': {'html': f'<p>SCORM Course: {title}</p>', 'markdown': title}}]
            })

        return {
            'title': title,
            'modules': [
                {
                    'title': 'Module 1: SCORM Package Items',
                    'summary': f"Version: SCORM {manifest_data.get('version')}",
                    'lessons': lessons
                }
            ]
        }
